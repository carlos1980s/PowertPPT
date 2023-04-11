import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Sample data with month and value columns
data = [
    {'month': '2021-01', 'value': 100},
    {'month': '2021-02', 'value': 120},
    {'month': '2021-03', 'value': 110},
    {'month': '2021-04', 'value': 130},
]

# Create a pandas DataFrame from the data
df = pd.DataFrame(data)

# Calculate the month-on-month percentage change
df['percentage_change'] = df['value'].pct_change() * 100

# Create a new presentation
presentation = Presentation()

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add a table to the slide
table_shape = slide.shapes.add_table(rows=len(df) + 1, cols=3, left=Inches(2), top=Inches(2), width=Inches(6), height=Inches(3))
table = table_shape.table

# Set table headers
table.cell(0, 0).text = 'Month'
table.cell(0, 1).text = 'Value'
table.cell(0, 2).text = 'Percentage Change'

# Fill the table with data
for index, row in df.iterrows():
    table.cell(index + 1, 0).text = row['month']
    table.cell(index + 1, 1).text = str(row['value'])

    # Calculate the percentage change and add delta symbol
    if not pd.isna(row['percentage_change']):
        value_cell = table.cell(index + 1, 2)
        run = value_cell.text_frame.paragraphs[0].add_run()
        run.text = 'Δ {:.2f}%'.format(row['percentage_change'])

        # Set the color based on the percentage change
        if row['percentage_change'] >= 0:
            run.font.color.rgb = RGBColor(0, 128, 0)  # Green
        else:
            run.font.color.rgb = RGBColor(255, 0, 0)  # Red

# Save the presentation
presentation.save('month_on_month_changes.pptx')
